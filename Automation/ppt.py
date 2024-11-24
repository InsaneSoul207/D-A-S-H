from pptx import Presentation
import google.generativeai as genai
import ast


def generate_response(query):
    genai.configure(api_key="AIzaSyAOR0bY6gixmV9iflLV6LkMnwExx4M2B1c")
    model = genai.GenerativeModel("gemini-1.5-flash")
    chat = model.start_chat(
        history=[
            {"role": "user", "parts": "Hello"},
            {"role": "model", "parts": "Hello! I am Alpha. I can surely not harm humans. HaHa. I am developed by Eshaan Mishra, the main function of Alpha is to make learning and education easier and more convenient for students. My purpose is to simplify your learning journey by providing personalized assistance, innovative teaching methods, and tailored resources to meet your unique needs. I am here to make your educational experience more enjoyable and effective. Feel free to ask me any questions or let me know how I can assist you in your learning adventure! and also in many more things from your life."},
        ]
    )
    response = chat.send_message(query)
    return response.text

def create_presentation(slides):
    prs = Presentation()

    for slide_info in slides:
        slide_number, title, content = slide_info
        
        # Create a slide with the title and content
        slide_layout = prs.slide_layouts[1]  # Using the title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]  # Assuming the second placeholder is for content

        title_shape.text = title
        content_shape.text = content

    return prs


if __name__ == "__main__":
    text = generate_response(f"""create a ppt on AI revolution
                            in the following format:
                            this should be a python list of tupples where each tupple is in the format:(slide_number, title, content)
                            i dont want a single letter more than this
                            """)
    text = text.replace("python","")
    text = text.replace("`","")
    text = ast.literal_eval(text)           
    presentation = create_presentation(text)
    ppt = input("Enter: ")
    presentation.save(f"{ppt}.pptx")